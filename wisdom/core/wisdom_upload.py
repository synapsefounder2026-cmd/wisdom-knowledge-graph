"""
Wisdom Manual Upload Pipeline
Ho tro: PDF, Word, PPT, Excel, TXT, MD, Audio, Video, Image, EPUB

Fixes:
  P-004: save_to_neo4j returns (file_id, neo4j_node_id) — real Neo4j elementId
  P-012: SHA-256 dedup check truoc khi ingest
  RULE-B: Them valid_from, source_type fields

Usage:
    python wisdom_upload.py <file_path>
"""

import os
import sys
import hashlib
import json
import re
import requests
from datetime import datetime
from pathlib import Path

try:
    import sys as _sys, os as _os
    _sys.path.append(_os.path.dirname(_os.path.abspath(__file__)))
    from wisdom_dedup import WisdomDedup
    _dedup = WisdomDedup()
except Exception as e:
    print(f"  Dedup warning: {e}")
    _dedup = None

# ── Strip emoji ───────────────────────────────────────────────────────────────
def strip_emoji(text: str) -> str:
    if not isinstance(text, str):
        return str(text) if text else ""
    emoji_pattern = re.compile(
        "["
        u"\U0001F600-\U0001F64F"
        u"\U0001F300-\U0001F5FF"
        u"\U0001F680-\U0001F6FF"
        u"\U0001F1E0-\U0001F1FF"
        u"\U00002600-\U000027BF"
        u"\U0001F900-\U0001F9FF"
        "]+",
        flags=re.UNICODE,
    )
    return emoji_pattern.sub("", text).strip()

# ── Optional imports ──────────────────────────────────────────────────────────
try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except ImportError:
    HAS_OCR = False

try:
    import ebooklib
    from ebooklib import epub
    from html.parser import HTMLParser
    HAS_EPUB = True
except ImportError:
    HAS_EPUB = False

from neo4j import GraphDatabase
from qdrant_client import QdrantClient
from qdrant_client.models import PointStruct, VectorParams, Distance

# ── Config ────────────────────────────────────────────────────────────────────
OLLAMA_BASE  = "http://localhost:11434"
OLLAMA_MODEL = "llama3.1:8b"
EMBED_MODEL  = "nomic-embed-text"
NEO4J_URI    = "bolt://localhost:7687"
NEO4J_USER   = "neo4j"
NEO4J_PASS   = "password123"
QDRANT_HOST  = "localhost"
QDRANT_PORT  = 6333
COLLECTION   = "wisdom_knowledge"
VECTOR_SIZE  = 768

SUPPORTED = {
    "document": [".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".txt", ".md"],
    "audio":    [".mp3", ".wav", ".m4a", ".ogg", ".flac"],
    "video":    [".mp4", ".mov", ".avi", ".mkv", ".webm"],
    "image":    [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"],
    "ebook":    [".epub"],
}

# ── Text Extractors ───────────────────────────────────────────────────────────

def extract_pdf(path: str) -> str:
    if not HAS_PDF:
        return "[PyPDF2 not installed]"
    reader = PdfReader(path)
    text = [t for page in reader.pages for t in [page.extract_text()] if t]
    if text and len(" ".join(text).strip()) > 50:
        return "\n".join(text)
    print("  Text-based failed, trying OCR...")
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(path, poppler_path=r"C:\poppler\poppler-25.12.0\Library\bin")
        ocr_text = []
        for i, img in enumerate(images):
            print(f"  OCR page {i+1}/{len(images)}...")
            t = pytesseract.image_to_string(img, lang="eng")
            if t.strip():
                ocr_text.append(t)
        return "\n".join(ocr_text)
    except Exception as e:
        return f"[OCR failed: {e}]"


def extract_docx(path: str) -> str:
    if not HAS_DOCX:
        return "[python-docx not installed]"
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_pptx(path: str) -> str:
    if not HAS_PPTX:
        return "[python-pptx not installed]"
    prs = Presentation(path)
    text = []
    for i, slide in enumerate(prs.slides):
        parts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if parts:
            text.append(f"[Slide {i+1}] " + " | ".join(parts))
    return "\n".join(text)


def extract_xlsx(path: str) -> str:
    if not HAS_XLSX:
        return "[openpyxl not installed]"
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(c) for c in row if c is not None])
            if row_text.strip():
                text.append(row_text)
    return "\n".join(text)


def extract_txt(path: str) -> str:
    for enc in ["utf-8", "utf-16", "latin-1"]:
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except Exception:
            continue
    return "[Could not read text file]"


def extract_image_ocr(path: str) -> str:
    if not HAS_OCR:
        return "[pytesseract/Pillow not installed]"
    try:
        return pytesseract.image_to_string(Image.open(path), lang="eng+vie")
    except Exception as e:
        return f"[OCR failed: {e}]"


def extract_epub(path: str) -> str:
    if not HAS_EPUB:
        return "[ebooklib not installed]"

    class MLStripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.fed = []
        def handle_data(self, d):
            self.fed.append(d)
        def get_data(self):
            return " ".join(self.fed)

    book = epub.read_epub(path)
    text = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            s = MLStripper()
            s.feed(item.get_content().decode("utf-8", errors="ignore"))
            t = s.get_data().strip()
            if t:
                text.append(t)
    return "\n".join(text)


def transcribe_audio(path: str) -> str:
    print("  Transcribing audio via Groq...")
    groq_key = os.environ.get("GROQ_API_KEY", "")
    if not groq_key:
        env_path = os.path.expanduser("~/.config/watch-cli/env")
        if os.path.exists(env_path):
            with open(env_path, encoding="utf-8") as f:
                for line in f:
                    if line.startswith("GROQ_API_KEY="):
                        groq_key = line.split("=", 1)[1].strip()
                        break
    if not groq_key:
        return "[No GROQ_API_KEY found]"
    with open(path, "rb") as f:
        response = requests.post(
            "https://api.groq.com/openai/v1/audio/transcriptions",
            headers={"Authorization": f"Bearer {groq_key}"},
            files={"file": (os.path.basename(path), f)},
            data={"model": "whisper-large-v3-turbo", "response_format": "text"},
            timeout=120,
        )
    return response.text.strip() if response.status_code == 200 else f"[Transcription failed: {response.text}]"


def transcribe_video(path: str) -> str:
    import subprocess
    print("  Extracting audio from video...")
    audio_path = path.rsplit(".", 1)[0] + "_audio.mp3"
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", path, "-vn", "-ac", "1", "-ar", "16000", "-b:a", "32k", "-f", "mp3", audio_path],
            capture_output=True, check=True,
        )
        text = transcribe_audio(audio_path)
        os.remove(audio_path)
        return text
    except Exception as e:
        return f"[Video transcription failed: {e}]"


def detect_type(path: str) -> tuple[str, str]:
    ext = Path(path).suffix.lower()
    for category, exts in SUPPORTED.items():
        if ext in exts:
            return category, ext
    return "unknown", ext


def extract_content(path: str) -> str:
    category, ext = detect_type(path)
    print(f"  Extracting content [{category}|{ext}]...")
    if ext == ".pdf":               return extract_pdf(path)
    elif ext in [".docx", ".doc"]:  return extract_docx(path)
    elif ext in [".pptx", ".ppt"]:  return extract_pptx(path)
    elif ext in [".xlsx", ".xls"]:  return extract_xlsx(path)
    elif ext in [".txt", ".md"]:    return extract_txt(path)
    elif ext == ".epub":            return extract_epub(path)
    elif category == "image":       return extract_image_ocr(path)
    elif category == "audio":       return transcribe_audio(path)
    elif category == "video":       return transcribe_video(path)
    else:                           return f"[Unsupported format: {ext}]"

# ── AI Analysis ───────────────────────────────────────────────────────────────

def analyze_with_ollama(content: str, filename: str) -> dict:
    content  = strip_emoji(content)
    filename = strip_emoji(filename)
    print(f"  Analyzing with {OLLAMA_MODEL}...")
    prompt = f"""Analyze this document and extract structured knowledge.
Return ONLY valid JSON, no markdown, no explanation.

Filename: {filename}
Content (first 3000 chars): {content[:3000]}

Return this exact JSON structure:
{{
  "title": "document title or topic",
  "summary": "2-3 sentence summary",
  "key_concepts": ["concept1", "concept2", "concept3"],
  "insights": ["insight1", "insight2"],
  "tags": ["tag1", "tag2", "tag3"],
  "language": "vi or en",
  "document_type": "book/article/report/lecture/note/data/other",
  "value_flywheel": "learning/experience/earning/contribution/growth"
}}"""

    response = requests.post(
        f"{OLLAMA_BASE}/api/generate",
        json={"model": OLLAMA_MODEL, "prompt": prompt, "stream": False},
        timeout=120,
    )
    raw = response.json().get("response", "{}").strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    try:
        return json.loads(raw.strip())
    except Exception:
        return {
            "title": filename,
            "summary": content[:200],
            "key_concepts": [],
            "insights": [],
            "tags": [],
            "language": "en",
            "document_type": "other",
            "value_flywheel": "learning",
        }

# ── Storage ───────────────────────────────────────────────────────────────────

def get_embedding(text: str) -> list:
    response = requests.post(
        f"{OLLAMA_BASE}/api/embeddings",
        json={"model": EMBED_MODEL, "prompt": text},
        timeout=60,
    )
    return response.json().get("embedding", [])


def save_to_neo4j(file_id: str, path: str, analysis: dict) -> tuple[str, object]:
    """
    Luu Document node vao Neo4j.

    Returns:
        (file_id, neo4j_node_id)
        neo4j_node_id: elementId(d) — bridge sang Qdrant (P-004)
    """
    print("  Saving to Neo4j...")
    neo4j_node_id = None
    try:
        driver   = GraphDatabase.driver(NEO4J_URI, auth=(NEO4J_USER, NEO4J_PASS))
        filename = Path(path).name
        with driver.session() as session:
            # P-004: RETURN elementId(d)
            # RULE-B: du fields — them valid_from, source_type
            result = session.run(
                """
                MERGE (d:Document {id: $id})
                SET d.filename         = $filename,
                    d.path             = $path,
                    d.title            = $title,
                    d.summary          = $summary,
                    d.document_type    = $doc_type,
                    d.language         = $language,
                    d.value_flywheel   = $flywheel,
                    d.ingested_at      = $ingested_at,
                    d.trust_score      = 0.8,
                    d.decay_lambda     = 0.003,
                    d.valid_from       = $valid_from,
                    d.valid_until      = null,
                    d.epistemic_status = 'PENDING',
                    d.cultural_context = 'GLOBAL',
                    d.source_type      = 'DOCUMENT'
                RETURN elementId(d) AS node_id
                """,
                id=file_id,
                filename=strip_emoji(filename),
                path=path,
                title=strip_emoji(analysis.get("title", "")),
                summary=strip_emoji(analysis.get("summary", "")),
                doc_type=analysis.get("document_type", "other"),
                language=analysis.get("language", "en"),
                flywheel=analysis.get("value_flywheel", "learning"),
                ingested_at=datetime.now().isoformat(),
                valid_from=datetime.now().isoformat(),
            )
            record = result.single()
            if record:
                neo4j_node_id = record["node_id"]

            for concept in analysis.get("key_concepts", []):
                session.run(
                    """
                    MERGE (c:Concept {name: $name})
                    WITH c MATCH (d:Document {id: $doc_id})
                    MERGE (d)-[:HAS_CONCEPT]->(c)
                    """,
                    name=strip_emoji(concept), doc_id=file_id,
                )
            for tag in analysis.get("tags", []):
                session.run(
                    """
                    MERGE (t:Tag {name: $name})
                    WITH t MATCH (d:Document {id: $doc_id})
                    MERGE (d)-[:HAS_TAG]->(t)
                    """,
                    name=strip_emoji(tag), doc_id=file_id,
                )

        driver.close()
        print(f"  Neo4j: node_id={neo4j_node_id} | {len(analysis.get('key_concepts', []))} concepts saved")
    except Exception as e:
        print(f"  Neo4j ERROR: {e}")
    return file_id, neo4j_node_id


def save_to_qdrant(file_id: str, neo4j_node_id, path: str, content: str, analysis: dict):
    """P-004: nhan neo4j_node_id lam payload bridge."""
    print("  Saving to Qdrant...")
    try:
        client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)
        existing = [c.name for c in client.get_collections().collections]
        if COLLECTION not in existing:
            client.create_collection(
                collection_name=COLLECTION,
                vectors_config=VectorParams(size=VECTOR_SIZE, distance=Distance.COSINE),
            )
        text_to_embed = f"{analysis.get('title', '')} {analysis.get('summary', '')} {content[:1000]}"
        embedding = get_embedding(text_to_embed)
        if len(embedding) != VECTOR_SIZE:
            print(f"  WARNING: Embedding size mismatch, skipping Qdrant")
            return
        point_id = int(hashlib.md5(file_id.encode()).hexdigest()[:8], 16)
        client.upsert(
            collection_name=COLLECTION,
            points=[PointStruct(
                id=point_id,
                vector=embedding,
                payload={
                    "neo4j_node_id": neo4j_node_id,  # P-004
                    "file_id":       file_id,
                    "filename":      Path(path).name,
                    "path":          path,
                    "title":         analysis.get("title", ""),
                    "summary":       analysis.get("summary", ""),
                    "tags":          analysis.get("tags", []),
                    "key_concepts":  analysis.get("key_concepts", []),
                    "document_type": analysis.get("document_type", "other"),
                    "value_flywheel": analysis.get("value_flywheel", "learning"),
                    "source":        "manual_upload",
                    "ingested_at":   datetime.now().isoformat(),
                },
            )]
        )
        print(f"  Qdrant: vector saved | neo4j_node_id={neo4j_node_id}")
    except Exception as e:
        print(f"  Qdrant ERROR: {e}")

# ── Main ──────────────────────────────────────────────────────────────────────

def upload(path: str):
    path = os.path.abspath(path)
    print(f"\n{'='*60}")
    print("  WISDOM MANUAL UPLOAD")
    print(f"{'='*60}")
    print(f"  File: {Path(path).name}")

    if not os.path.exists(path):
        print(f"  File not found: {path}")
        return None

    category, ext = detect_type(path)
    if category == "unknown":
        print(f"  Unsupported format: {ext}")
        return None

    print(f"  Type: {category} ({ext})")
    print(f"  Size: {os.path.getsize(path) / 1024:.1f} KB\n")

    # P-012: Dedup check tren file hash
    file_hash = hashlib.sha256(open(path, "rb").read()).hexdigest()
    if _dedup is not None:
        try:
            if _dedup.check_duplicate(file_hash).get("is_duplicate", False):
                print(f"  [DEDUP] File da ton tai, skip: {Path(path).name}")
                return None
        except Exception as e:
            print(f"  [DEDUP] Check failed (non-blocking): {e}")

    # Step 1: Extract
    content = extract_content(path)
    if not content or len(content.strip()) < 10:
        print("  Could not extract content from file.")
        return None
    print(f"  Extracted: {len(content)} chars")

    # P-012: Dedup check tren content hash
    content_hash = hashlib.sha256(content.encode()).hexdigest()
    if _dedup is not None:
        try:
            if _dedup.check_duplicate(content_hash).get("is_duplicate", False):
                print("  [DEDUP] Noi dung da ton tai (content match), skip.")
                return None
        except Exception as e:
            print(f"  [DEDUP] Content check failed (non-blocking): {e}")

    # Step 2: Analyze
    filename = Path(path).name
    analysis = analyze_with_ollama(content, filename)
    print(f"  Title: {analysis.get('title')}")
    print(f"  Concepts: {', '.join(analysis.get('key_concepts', []))}")

    # Step 3: Generate file_id
    file_id = hashlib.md5(f"{path}{os.path.getmtime(path)}".encode()).hexdigest()[:12]

    # Step 4: Neo4j — P-004: lay neo4j_node_id
    file_id, neo4j_node_id = save_to_neo4j(file_id, path, analysis)

    # Step 5: Qdrant — P-004: truyen neo4j_node_id
    save_to_qdrant(file_id, neo4j_node_id, path, content, analysis)

    # P-012: Register hash sau khi thanh cong
    if _dedup is not None:
        try:
            _dedup.register_checksum(str(neo4j_node_id), file_hash, path, datetime.now().isoformat())
            _dedup.register_checksum(str(neo4j_node_id), content_hash, path, datetime.now().isoformat())
        except Exception as e:
            print(f"  [DEDUP] Register failed (non-blocking): {e}")

    print(f"\n{'='*60}")
    print("  UPLOAD COMPLETE")
    print(f"  file_id:       {file_id}")
    print(f"  neo4j_node_id: {neo4j_node_id}")
    print(f"  Searchable via: python wisdom_query.py \"<question>\"")
    print(f"{'='*60}\n")

    return {
        "file_id":       file_id,
        "neo4j_node_id": neo4j_node_id,
        "analysis":      analysis,
    }


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python wisdom_upload.py <file_path>")
        sys.exit(1)
    upload(sys.argv[1])
